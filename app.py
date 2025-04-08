from flask import Flask, request, jsonify
from flask_cors import CORS
from docx import Document as DocxDocument
from datetime import datetime
from openai import OpenAI
import os
import tempfile
import docx2txt
import pandas as pd
import pdfplumber
import pptx
from PIL import Image
import pytesseract
import time

app = Flask(__name__)
CORS(app, origins=["https://financial-optimization-dashboard.netlify.app"])

REPORT_FOLDER = os.path.join(app.root_path, 'static', 'reports')
os.makedirs(REPORT_FOLDER, exist_ok=True)

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

def extract_text(file_storage):
    filename = file_storage.filename.lower()
    try:
        if filename.endswith(".pdf"):
            with pdfplumber.open(file_storage.stream) as pdf:
                return "\n".join(page.extract_text() or "" for page in pdf.pages)
        elif filename.endswith(".docx"):
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                file_storage.save(tmp.name)
                return docx2txt.process(tmp.name)
        elif filename.endswith(".pptx"):
            prs = pptx.Presentation(file_storage)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif filename.endswith(".png") or filename.endswith(".jpg") or filename.endswith(".jpeg"):
            image = Image.open(file_storage.stream)
            return pytesseract.image_to_string(image)
        elif filename.endswith(".xlsx"):
            xls = pd.read_excel(file_storage, sheet_name=None)
            return "\n".join(df.to_string() for df in xls.values())
        elif filename.endswith(".csv"):
            df = pd.read_csv(file_storage)
            return df.to_string()
        else:
            return file_storage.read().decode("utf-8", errors="ignore")
    except Exception as e:
        return f"[Error reading {filename}: {e}]"

def trim_text(text, max_chars=6000):
    return text[:max_chars]

def generate_section(title, instruction, context):
    trimmed_context = trim_text(context)
    messages = [
        {"role": "system", "content": "You are a financial advisor generating professional financial optimization reports."},
        {"role": "user", "content": f"{instruction}\n\nBusiness Context:\n{trimmed_context}"}
    ]
    for _ in range(3):  # Retry up to 3 times
        try:
            response = client.chat.completions.create(
                model="gpt-4",
                messages=messages,
                temperature=0.7
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            if "rate limit" in str(e).lower() or "429" in str(e):
                time.sleep(5)
            else:
                return f"Error generating this section: {e}"
    return "Error: GPT failed after multiple attempts."

def add_example_table(doc, section_title):
    if "Expense Breakdown" in section_title:
        table = doc.add_table(rows=1, cols=3)
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Expense Category'
        hdr_cells[1].text = 'Monthly Total'
        hdr_cells[2].text = 'Percentage'
        rows = [
            ('Payroll', '$12,000', '48%'),
            ('Software Subscriptions', '$3,000', '12%'),
            ('Utilities', '$1,000', '4%'),
            ('Marketing', '$5,000', '20%'),
            ('Other', '$4,000', '16%')
        ]
        for row in rows:
            row_cells = table.add_row().cells
            for i, val in enumerate(row):
                row_cells[i].text = val
        doc.add_paragraph("")  # Spacer
    elif "Profitability Analysis" in section_title:
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Table Grid'
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Metric'
        hdr_cells[1].text = 'Q1'
        hdr_cells[2].text = 'Q2'
        hdr_cells[3].text = 'Change'
        rows = [
            ('Revenue', '$120,000', '$130,000', '+8.3%'),
            ('Gross Profit', '$45,000', '$52,000', '+15.6%'),
            ('Net Profit', '$8,000', '$11,500', '+43.8%')
        ]
        for row in rows:
            row_cells = table.add_row().cells
            for i, val in enumerate(row):
                row_cells[i].text = val
        doc.add_paragraph("")  # Spacer

@app.route('/')
def home():
    return "Financial Optimization Backend is Live!"

@app.route('/generate', methods=['POST'])
def generate_report():
    files = [request.files.get('file1'), request.files.get('file2'), request.files.get('file3')]
    context = ""

    for file in files:
        if file:
            context += extract_text(file) + "\n"

    if not context.strip():
        return jsonify({"error": "No valid file content found."}), 400

    doc = DocxDocument()
    doc.add_heading("Financial Metric Optimization Report", 0)

    sections = [
        ("Executive Summary", "Summarize key findings from financial documents and high-level trends."),
        ("1. Revenue & Income Patterns", "Analyze trends in sales, income sources, and pricing structures."),
        ("2. Expense Breakdown & Cost Management", "Assess fixed vs. variable costs and highlight opportunities to reduce expenses."),
        ("3. Profitability Analysis", "Evaluate gross margin, net margin, and profitability over time."),
        ("4. Cash Flow Insights", "Analyze inflows and outflows of cash to assess liquidity and runway."),
        ("5. Forecasting & Financial Risk", "Identify risk areas and provide projection insights."),
        ("6. Recommendations & Financial Optimization", "Offer specific strategies to improve financial health."),
        ("Conclusion", "Wrap up the financial overview and suggest next steps for improvement.")
    ]

    for title, instruction in sections:
        doc.add_heading(title, level=1)
        section_text = generate_section(title, instruction, context)
        doc.add_paragraph(section_text)
        add_example_table(doc, title)

    filename = f"financial_report_{datetime.now().strftime('%Y%m%d%H%M%S')}.docx"
    file_path = os.path.join(REPORT_FOLDER, filename)
    doc.save(file_path)

    return jsonify({'download_url': f'/static/reports/{filename}'})

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)